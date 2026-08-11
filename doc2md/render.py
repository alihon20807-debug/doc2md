"""Rasterize input documents (PDFs, PPT/PPTX presentations, or standalone images) into per-page PIL images."""

from __future__ import annotations

import asyncio
import io
import os
import shutil
import subprocess
import tempfile
from concurrent.futures import ProcessPoolExecutor
from pathlib import Path

import pymupdf as fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont

from doc2md.models import PageImage

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}
PPT_SUFFIXES = {".ppt", ".pptx"}


def _render_pdf_page(args: tuple[str, int, int, str]) -> PageImage:
    pdf_path_str, page_index, dpi, source_name = args
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    with fitz.open(pdf_path_str) as doc:
        page = doc.load_page(page_index)
        pix = page.get_pixmap(matrix=matrix, colorspace=fitz.csRGB, alpha=False)
        image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        return PageImage(page_no=page_index + 1, image=image, source_name=source_name)


def render_pdf(path: Path, dpi: int) -> list[PageImage]:
    """Rasterize PDF pages in parallel using ProcessPoolExecutor."""
    with fitz.open(path) as doc:
        page_count = doc.page_count

    if page_count == 0:
        return []

    if page_count == 1:
        zoom = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)
        with fitz.open(path) as doc:
            page = doc.load_page(0)
            pix = page.get_pixmap(matrix=matrix, colorspace=fitz.csRGB, alpha=False)
            image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            return [PageImage(page_no=1, image=image, source_name=path.stem)]

    tasks = [(str(path.resolve()), i, dpi, path.stem) for i in range(page_count)]
    workers = min(page_count, os.cpu_count() or 4)
    with ProcessPoolExecutor(max_workers=workers) as executor:
        pages = list(executor.map(_render_pdf_page, tasks))
    return sorted(pages, key=lambda p: p.page_no)


def render_ppt(path: Path, dpi: int) -> list[PageImage]:
    """Convert PPT/PPTX presentation to page images."""
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir)
        converted_pdf = None

        # 1. Try LibreOffice / soffice CLI
        for cmd in ["soffice", "libreoffice"]:
            if shutil.which(cmd):
                try:
                    res = subprocess.run(
                        [cmd, "--headless", "--convert-to", "pdf", "--outdir", str(tmp_path), str(path.resolve())],
                        capture_output=True,
                        timeout=120,
                    )
                    if res.returncode == 0:
                        candidate = tmp_path / f"{path.stem}.pdf"
                        if candidate.exists():
                            converted_pdf = candidate
                            break
                except Exception:
                    pass

        # 2. Try Windows PowerPoint COM if on Windows and soffice failed
        if converted_pdf is None and os.name == "nt":
            try:
                import win32com.client  # type: ignore

                ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                pdf_target = tmp_path / f"{path.stem}.pdf"
                deck = ppt_app.Presentations.Open(str(path.resolve()), WithWindow=False)
                deck.SaveAs(str(pdf_target.resolve()), 32)  # 32 = ppSaveAsPDF
                deck.Close()
                if pdf_target.exists():
                    converted_pdf = pdf_target
            except Exception:
                pass

        if converted_pdf is not None:
            return render_pdf(converted_pdf, dpi=dpi)

        # 3. Fallback: python-pptx slide rendering / image extraction
        try:
            import pptx
        except ImportError:
            raise RuntimeError(
                f"Cannot render PPT presentation '{path}'. Please install LibreOffice, Microsoft PowerPoint, "
                "or python-pptx."
            )

        prs = pptx.Presentation(str(path))
        pages: list[PageImage] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_w = int((prs.slide_width or 9144000) / 9144000 * (dpi * 10))
            slide_h = int((prs.slide_height or 5143500) / 9144000 * (dpi * 10))
            canvas_img = Image.new("RGB", (max(800, slide_w), max(600, slide_h)), "white")
            draw = ImageDraw.Draw(canvas_img)
            font = ImageFont.load_default()

            y_offset = 40
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    draw.text((40, y_offset), shape.text.strip(), fill="black", font=font)
                    y_offset += 30
                elif hasattr(shape, "image") and getattr(shape.image, "blob", None):
                    try:
                        shape_img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                        canvas_img.paste(shape_img, (40, y_offset))
                        y_offset += shape_img.height + 20
                    except Exception:
                        pass
            pages.append(PageImage(page_no=idx, image=canvas_img, source_name=path.stem))
        return pages


def load_image_files(paths: list[Path], source_name: str) -> list[PageImage]:
    pages: list[PageImage] = []
    for page_index, image_path in enumerate(sorted(paths), start=1):
        image = Image.open(image_path).convert("RGB")
        pages.append(PageImage(page_no=page_index, image=image, source_name=source_name))
    return pages


def render_input(path: Path, dpi: int) -> list[PageImage]:
    """Dispatch on the input: a PDF file, PPT/PPTX presentation, a single image, or a directory of page images."""
    if path.is_dir():
        image_paths = [p for p in path.iterdir() if p.suffix.lower() in IMAGE_SUFFIXES]
        if not image_paths:
            raise ValueError(f"No page images (extensions {sorted(IMAGE_SUFFIXES)}) found in {path}")
        return load_image_files(image_paths, source_name=path.name)

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return render_pdf(path, dpi=dpi)
    if suffix in PPT_SUFFIXES:
        return render_ppt(path, dpi=dpi)
    if suffix in IMAGE_SUFFIXES:
        return load_image_files([path], source_name=path.stem)

    raise ValueError(
        f"Unsupported input '{path}': expected a .pdf file, .ppt/.pptx presentation, "
        "an image file, or a directory of images"
    )


async def render_input_async(path: Path, dpi: int) -> list[PageImage]:
    """Async wrapper around render_input offloading rendering off the main asyncio loop."""
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, render_input, path, dpi)

